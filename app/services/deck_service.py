import asyncio
import re
import uuid
from io import BytesIO
from pathlib import Path

import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.config import settings
from app.models import DeckDraft, GenerationRequest


PALETTES = {
    1: {"bg": "0B1020", "panel": "121A31", "text": "F4F7FB", "muted": "A5B1C7", "accent": "68D5FF"},
    2: {"bg": "FFF8F1", "panel": "FFFFFF", "text": "102033", "muted": "617089", "accent": "FF7A29"},
    3: {"bg": "F4F8F2", "panel": "FFFFFF", "text": "14281D", "muted": "60756B", "accent": "33B36B"},
    4: {"bg": "161616", "panel": "20232C", "text": "F7F8FA", "muted": "B8BFCC", "accent": "FF5A67"},
    5: {"bg": "F6FAFF", "panel": "FFFFFF", "text": "14243B", "muted": "697A8C", "accent": "3C82F6"},
    6: {"bg": "171112", "panel": "241B1D", "text": "FAF7F2", "muted": "C7B8AA", "accent": "F0A85C"},
}


async def build_presentation(request: GenerationRequest, deck: DeckDraft) -> Path:
    settings.generated_dir.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    palette = PALETTES.get(request.design_number, PALETTES[2])
    _build_cover_slide(presentation, deck, palette, request)

    for index, slide in enumerate(deck.slides, start=1):
        _build_content_slide(presentation, slide.title, slide.content, index, len(deck.slides), palette)

    safe_name = re.sub(r"[^a-zA-Z0-9_-]+", "-", deck.title.lower()).strip("-") or "deckmint"
    file_path = settings.generated_dir / f"{safe_name}-{uuid.uuid4().hex[:8]}.pptx"
    await asyncio.to_thread(presentation.save, str(file_path))
    return file_path


def _build_cover_slide(presentation: Presentation, deck: DeckDraft, palette: dict[str, str], request: GenerationRequest) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _fill_slide(slide, palette["bg"])
    _add_orb(slide, Inches(9.6), Inches(0.7), Inches(2.4), palette["accent"], 0.18)
    _add_orb(slide, Inches(10.4), Inches(4.8), Inches(1.5), palette["accent"], 0.12)

    accent_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.7), Inches(1.8), Inches(0.28)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = _rgb(palette["accent"])
    accent_bar.line.fill.background()

    eyebrow = slide.shapes.add_textbox(Inches(0.9), Inches(0.73), Inches(3.2), Inches(0.2))
    eyebrow_frame = eyebrow.text_frame
    eyebrow_frame.text = f"{request.audience.title()} Plan"
    eyebrow_frame.paragraphs[0].font.size = Pt(12)
    eyebrow_frame.paragraphs[0].font.bold = True
    eyebrow_frame.paragraphs[0].font.color.rgb = _rgb(palette["bg"])

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(6.2), Inches(2.2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.text = deck.title
    title_frame.paragraphs[0].font.name = "Aptos Display"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = _rgb(palette["text"])

    subtitle_box = slide.shapes.add_textbox(Inches(0.85), Inches(3.65), Inches(5.8), Inches(1.2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_frame.text = deck.subtitle
    subtitle_frame.paragraphs[0].font.name = "Aptos"
    subtitle_frame.paragraphs[0].font.size = Pt(15)
    subtitle_frame.paragraphs[0].font.color.rgb = _rgb(palette["muted"])

    metric_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(1.5), Inches(4.1), Inches(4.9)
    )
    metric_card.fill.solid()
    metric_card.fill.fore_color.rgb = _rgb(palette["panel"])
    metric_card.line.color.rgb = _rgb(palette["accent"])
    metric_card.line.width = Pt(1.4)

    metrics = [
        ("Slides", str(len(deck.slides))),
        ("Tone", request.tone.title()),
        ("Export", "PPTX"),
        ("Model", "Subscription-first"),
    ]
    top = 1.95
    for label, value in metrics:
        label_box = slide.shapes.add_textbox(Inches(8.6), Inches(top), Inches(1.4), Inches(0.25))
        label_box.text_frame.text = label.upper()
        label_box.text_frame.paragraphs[0].font.size = Pt(9)
        label_box.text_frame.paragraphs[0].font.color.rgb = _rgb(palette["muted"])
        label_box.text_frame.paragraphs[0].font.bold = True

        value_box = slide.shapes.add_textbox(Inches(8.55), Inches(top + 0.22), Inches(2.7), Inches(0.45))
        value_box.text_frame.text = value
        value_box.text_frame.paragraphs[0].font.size = Pt(19)
        value_box.text_frame.paragraphs[0].font.bold = True
        value_box.text_frame.paragraphs[0].font.color.rgb = _rgb(palette["text"])
        top += 1.0

    footer = slide.shapes.add_textbox(Inches(0.85), Inches(6.7), Inches(6), Inches(0.3))
    footer.text_frame.text = "DeckMint generates polished, editable presentations built for real business use."
    footer.text_frame.paragraphs[0].font.size = Pt(10)
    footer.text_frame.paragraphs[0].font.color.rgb = _rgb(palette["muted"])


def _build_content_slide(
    presentation: Presentation,
    title: str,
    content: str,
    index: int,
    total: int,
    palette: dict[str, str],
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _fill_slide(slide, palette["bg"])

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.55), Inches(12.1), Inches(6.3)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = _rgb(palette["panel"])
    panel.line.fill.background()

    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.85), Inches(0.85), Inches(0.12), Inches(1.05))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = _rgb(palette["accent"])
    stripe.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.15), Inches(0.82), Inches(6.2), Inches(0.8))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.name = "Aptos Display"
    title_box.text_frame.paragraphs[0].font.size = Pt(24)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = _rgb(palette["text"])

    body_box = slide.shapes.add_textbox(Inches(1.18), Inches(1.85), Inches(6.0), Inches(3.6))
    body_frame = body_box.text_frame
    body_frame.word_wrap = True
    body_frame.text = content
    body_frame.paragraphs[0].font.name = "Aptos"
    body_frame.paragraphs[0].font.size = Pt(16)
    body_frame.paragraphs[0].font.color.rgb = _rgb(palette["text"])
    body_frame.paragraphs[0].line_spacing = 1.25

    quote_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.1), Inches(1.2), Inches(3.7), Inches(4.8)
    )
    quote_panel.fill.solid()
    quote_panel.fill.fore_color.rgb = _rgb(palette["bg"])
    quote_panel.line.color.rgb = _rgb(palette["accent"])
    quote_panel.line.width = Pt(1)

    quote_text = slide.shapes.add_textbox(Inches(8.45), Inches(1.6), Inches(3.0), Inches(2.3))
    quote_text.text_frame.word_wrap = True
    quote_text.text_frame.text = "Professional structure, tighter messaging, and cleaner visuals are what make generated decks usable in the real world."
    quote_text.text_frame.paragraphs[0].font.size = Pt(15)
    quote_text.text_frame.paragraphs[0].font.bold = True
    quote_text.text_frame.paragraphs[0].font.color.rgb = _rgb(palette["text"])

    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.45), Inches(4.45), Inches(2.4), Inches(0.45)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = _rgb(palette["accent"])
    badge.line.fill.background()

    badge_text = slide.shapes.add_textbox(Inches(8.72), Inches(4.53), Inches(1.9), Inches(0.22))
    badge_text.text_frame.text = f"Slide {index:02d}"
    badge_text.text_frame.paragraphs[0].font.size = Pt(10)
    badge_text.text_frame.paragraphs[0].font.bold = True
    badge_text.text_frame.paragraphs[0].font.color.rgb = _rgb(palette["bg"])
    badge_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    footer = slide.shapes.add_textbox(Inches(10.85), Inches(6.35), Inches(0.8), Inches(0.22))
    footer.text_frame.text = f"{index}/{total}"
    footer.text_frame.paragraphs[0].font.size = Pt(10)
    footer.text_frame.paragraphs[0].font.color.rgb = _rgb(palette["muted"])
    footer.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def _fill_slide(slide, color_hex: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color_hex)


def _add_orb(slide, left, top, size, color_hex: str, transparency: float) -> None:
    orb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, size, size)
    orb.fill.solid()
    orb.fill.fore_color.rgb = _rgb(color_hex)
    orb.fill.transparency = transparency
    orb.line.fill.background()


def _rgb(color_hex: str) -> RGBColor:
    clean = color_hex.replace("#", "")
    return RGBColor.from_string(clean)
