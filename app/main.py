from __future__ import annotations
import os
import re
import uuid
import random
import asyncio
from io import BytesIO
from pathlib import Path
from typing import Optional, Dict, Any, List, Union
from typing import Any, Dict, List, Optional
import requests
from pptx import Presentation
from pptx.util import Inches

from fastapi import FastAPI, UploadFile, File, HTTPException, Query
from fastapi.responses import ORJSONResponse, FileResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi import FastAPI, UploadFile, File, HTTPException, Request
from fastapi.responses import ORJSONResponse
from pydantic import BaseModel

from app.config import settings
from app.models import (
    IdentifyRequest, IdentifyResponse,
    PortionRequest, PortionResponse, PortionEstimate,
    NutrientsRequest, NutrientsResponse,
    AnalyzeTextRequest, AnalyzeResponse, PPTGenerateRequest, PPTGenerateResponse
)
from app.clients.gemini import GeminiClient
from app.clients.perplexity import PerplexityClient
from app.utils import extract_json_object, ModelJSONError

app = FastAPI(
    title="Food Nutrients API",
    version="1.0.0",
    default_response_class=ORJSONResponse,
)

gemini = GeminiClient(api_key=settings.gemini_api_key)
pplx = PerplexityClient(api_key=settings.pplx_api_key)

GENERATED_DIR = Path(os.getenv("GENERATED_DIR", "GeneratedPresentations"))
CACHE_DIR = Path(os.getenv("CACHE_DIR", "Cache"))
DESIGNS_DIR = Path(os.getenv("DESIGNS_DIR", "Designs"))

GENERATED_DIR.mkdir(exist_ok=True)
CACHE_DIR.mkdir(exist_ok=True)
DESIGNS_DIR.mkdir(exist_ok=True)


# ---------------- Global JSON error handler (prevents jq explosions) ----------------

@app.exception_handler(Exception)
async def all_exception_handler(request: Request, exc: Exception):
    # Always return JSON so jq never breaks
    return ORJSONResponse(
        status_code=500,
        content={"error": "internal_error", "message": str(exc)[:2000]},
    )

# ---------------- Prompts ----------------

IDENTIFY_SYSTEM = "Return ONLY valid JSON. No markdown. No extra text."
PORTION_SYSTEM = "Return ONLY valid JSON. No markdown. No extra text."
NUTRIENTS_SYSTEM = (
    "Return ONLY valid JSON (no markdown, no commentary). "
    "Do NOT include citations, search_results, links, or extra fields."
)

def identify_user_prompt_text(text: str, hints: Optional[List[str]]) -> str:
    return f"""
Return JSON:
{{
  "candidates":[{{"name":str,"confidence":0-1,"normalized_name":str,"cuisine":str|null,"is_packaged":bool|null,"notes":str|null}}],
  "chosen": {{"name":str,"confidence":0-1,"normalized_name":str,"cuisine":str|null,"is_packaged":bool|null,"notes":str|null}}
}}
User text: {text}
Hints: {hints or []}
Rules:
- candidates max 3
- chosen MUST be an object (not string)
""".strip()

def identify_user_prompt_image(hints: Optional[List[str]]) -> str:
    return f"""
Return JSON:
{{
  "candidates":[{{"name":str,"confidence":0-1,"normalized_name":str,"cuisine":str|null,"is_packaged":bool|null,"notes":str|null}}],
  "chosen": {{"name":str,"confidence":0-1,"normalized_name":str,"cuisine":str|null,"is_packaged":bool|null,"notes":str|null}}
}}
Hints: {hints or []}
Rules:
- candidates max 3
- chosen MUST be an object (not string)
""".strip()

def portion_prompt_text(food_name: str, ctx: Optional[str]) -> str:
    return f"""
Food: {food_name}
Context: {ctx or ""}

Return ONLY JSON:
{{
  "servings": <float>,
  "grams_total": <float>,
  "items_count": <float|null>,
  "household": <string|null>,
  "confidence": <0-1>,
  "assumptions": ["max 3 items, each under 80 chars"]
}}
""".strip()

def portion_prompt_image(food_name: str, ctx: Optional[str]) -> str:
    return f"""
Food: {food_name}
Context: {ctx or ""}

Return ONLY JSON:
{{
  "servings": <float>,
  "grams_total": <float>,
  "items_count": <float|null>,
  "household": <string|null>,
  "confidence": <0-1>,
  "assumptions": ["max 3 items, each under 80 chars"]
}}
""".strip()

def nutrients_prompt(req: NutrientsRequest) -> str:
    p = req.portion
    return f"""
Return ONLY JSON with this exact schema keys:
food_name, portion, calories_kcal, macros, micros, vitamins, minerals, ingredients_guess, allergens_guess, data_sources, notes

Food: {req.food_name}
Region: {req.region}
Brand: {req.brand or ""}

Portion:
servings={p.servings}
grams_total={p.grams_total}
items_count={p.items_count}
household={p.household}

Rules:
- Use units g, mg, kcal
- macros/micros/vitamins/minerals values MUST be objects with:
  {{name, amount, unit, per_100g_amount, daily_value_percent}}
- ingredients_guess/allergens_guess/notes/data_sources MUST be arrays (lists)
- No citations, no links, no search_results.
""".strip()

PROMPT_TEMPLATE = """Write a presentation/powerpoint about the user's topic. You only answer with the presentation. Follow the structure of the example.

Notice:
- You do all the presentation text for the user.
- You write the texts no longer than 250 characters!
- You make very short titles!
- You make the presentation easy to understand.
- The presentation has a table of contents.
- The presentation has a summary.
- At least 7 slides.
- For each slide, after the #Content: line, add an #Image: line describing a relevant image that could visually represent the slide's topic.
- If no image is relevant, write #Image: none.

Example! - Stick to this formatting exactly!
#Title: TITLE OF THE PRESENTATION

#Slide: 1
#Header: table of contents
#Content: 1. CONTENT OF THIS POWERPOINT
2. CONTENTS OF THIS POWERPOINT
3. CONTENT OF THIS POWERPOINT
#Image: a 3D illustration of a table of contents in a book

#Slide: 2
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE
#Image: relevant illustration description here

#Slide: 3
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE
#Image: relevant illustration description here

#Slide: 4
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE
#Image: relevant illustration description here

#Slide: 5
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE
#Image: relevant illustration description here

#Slide: 6
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE
#Image: relevant illustration description here

#Slide: 7
#Header: summary
#Content: CONTENT OF THE SUMMARY
#Image: an illustration of a person reviewing a summary report

#Slide: END"""

def _local_fallback_presentation(topic: str) -> str:
    t = (topic or "Your Topic").title()
    return f"""#Title: {t}

#Slide: 1
#Header: table of contents
#Content: 1. Intro
2. Why it matters
3. Key points
4. Examples
5. Tips
6. Pitfalls
7. Summary
#Image: a 3D illustration of a table of contents in a book

#Slide: 2
#Header: intro
#Content: Brief overview of {topic}. Scope and goal.
#Image: relevant illustration

#Slide: 3
#Header: why it matters
#Content: Impact, use-cases, and benefits in simple terms.
#Image: simple infographic

#Slide: 4
#Header: key points
#Content: 3–5 core ideas about {topic}.
#Image: icons representing key ideas

#Slide: 5
#Header: examples
#Content: 2–3 quick examples.
#Image: storyboard-like illustration

#Slide: 6
#Header: tips & pitfalls
#Content: Do's and don'ts.
#Image: checklist illustration

#Slide: 7
#Header: summary
#Content: Short recap and next steps.
#Image: an illustration of a person reviewing a summary report

#Slide: END
"""

async def generate_presentation_text_async(topic: str) -> str:
    # IMPORTANT: uses YOUR PerplexityClient (same service instance)
    system = PROMPT_TEMPLATE
    user = f"The user wants a presentation about {topic}"

    try:
        content = await pplx.chat(
            model=settings.pplx_sonar_model,
            system=system,
            user=user,
            temperature=0.35,
            max_tokens=1400,
            search_recency_filter="month",
        )
        if "#Title:" not in content or "#Slide:" not in content:
            return ""
        return content
    except Exception:
        return ""

async def create_presentation_async(text_content: str, design_number: int, presentation_name: str) -> str:
    template_path = DESIGNS_DIR / f"Design-{design_number}.pptx"
    if not template_path.exists():
        template_path = DESIGNS_DIR / "Design-1.pptx"

    if template_path.exists():
        try:
            prs = Presentation(str(template_path))
        except Exception:
            prs = Presentation()
    else:
        prs = Presentation()

    # Layout helpers
    def _safe_layout(idx: int) -> int:
        if len(prs.slide_layouts) == 0:
            return 0
        return max(0, min(idx, len(prs.slide_layouts) - 1))

    slide_layout_index = _safe_layout(1 if len(prs.slide_layouts) > 1 else 0)
    slide_placeholder_index = 1 if len(prs.slide_layouts) > 1 else 0

    slide_title = ""
    slide_content = ""
    slide_image_prompt = None
    slide_count = 0
    last_layout = -1
    first_time = True

    lines = [ln.rstrip("\n") for ln in text_content.splitlines()]
    i = 0

    async def commit_slide():
        nonlocal slide_title, slide_content, slide_image_prompt, slide_count
        if slide_count > 0 and (slide_title or slide_content):
            slide = prs.slides.add_slide(prs.slide_layouts[_safe_layout(slide_layout_index)])

            # title
            try:
                slide.shapes.title.text = slide_title
            except Exception:
                try:
                    slide.placeholders[0].text = slide_title
                except Exception:
                    pass

            # body
            try:
                body = slide.shapes.placeholders[slide_placeholder_index]
                if hasattr(body, "text_frame"):
                    body.text_frame.text = slide_content
                else:
                    body.text = slide_content
            except Exception:
                pass

            # NOTE: Image generation removed for Render stability (g4f is flaky & heavy).
            # You can add it back later behind a feature flag.

    while i < len(lines):
        line = lines[i]

        if line.startswith("#Title:"):
            title = line.replace("#Title:", "").strip()
            slide = prs.slides.add_slide(prs.slide_layouts[_safe_layout(0)])
            try:
                slide.shapes.title.text = title
            except Exception:
                pass
            i += 1
            continue

        if line.startswith("#Slide:"):
            await commit_slide()
            slide_count += 1
            slide_title = ""
            slide_content = ""
            slide_image_prompt = None

            # rotate layouts a bit
            layouts = [1, 7, 8] if len(prs.slide_layouts) >= 9 else ([1] if len(prs.slide_layouts) > 1 else [0])
            if first_time:
                slide_layout_index = _safe_layout(1 if len(prs.slide_layouts) > 1 else 0)
                slide_placeholder_index = 1 if len(prs.slide_layouts) > 1 else 0
                first_time = False
            else:
                nxt = last_layout
                tries = 0
                while nxt == last_layout and tries < 10:
                    nxt = random.choice(layouts)
                    tries += 1
                slide_layout_index = _safe_layout(nxt)
                slide_placeholder_index = 2 if slide_layout_index == 8 else 1

            last_layout = slide_layout_index
            i += 1
            continue

        if line.startswith("#Header:"):
            slide_title = line.replace("#Header:", "").strip()
            i += 1
            continue

        if line.startswith("#Content:"):
            slide_content = line.replace("#Content:", "").strip()
            i += 1
            while i < len(lines) and not lines[i].startswith("#"):
                slide_content += "\n" + lines[i]
                i += 1
            continue

        if line.startswith("#Image:"):
            slide_image_prompt = line.replace("#Image:", "").strip()
            i += 1
            continue

        i += 1

    await commit_slide()

    out_path = GENERATED_DIR / f"{presentation_name}.pptx"
    await asyncio.to_thread(prs.save, str(out_path))
    return str(out_path)


async def compute_nutrients(req: NutrientsRequest) -> NutrientsResponse:
    out = await pplx.chat(
        model=settings.pplx_sonar_model,
        system=NUTRIENTS_SYSTEM,
        user=nutrients_prompt(req),
        temperature=0.1,
        max_tokens=1400,
        disable_search=True,                      #  [oai_citation:5‡docs.perplexity.ai](https://docs.perplexity.ai/api-reference/chat-completions-post)
        response_format={"type": "json_object"},  #  [oai_citation:6‡docs.perplexity.ai](https://docs.perplexity.ai/api-reference/chat-completions-post)
    )

    # Since response_format is json_object, `out` should already be valid JSON.
    try:
        raw = json.loads(out)
        if not isinstance(raw, dict):
            raise ValueError("not an object")
    except Exception:
        # Last resort repair via Gemini
        raw = await force_json_with_gemini("NutrientsResponse schema JSON object", out)

    raw = normalize_nutrients_obj(raw, req)
    return NutrientsResponse.model_validate(raw)

# ---------------- Helpers ----------------

def model_json_or_400(text: str) -> Any:
    try:
        return extract_json_object(text)
    except ModelJSONError as e:
        raise HTTPException(status_code=400, detail=str(e))

async def force_json_with_gemini(schema_hint: str, text: str) -> Dict[str, Any]:
    repair_system = "You repair JSON. Return ONLY valid JSON. No markdown. No extra text."
    repair_prompt = f"""
Fix the following invalid/truncated JSON into valid JSON.

Schema hint:
{schema_hint}

Input:
{text}

Return ONLY JSON. If impossible, return {{}}.
""".strip()
    repaired = await gemini.generate_text(
        model=settings.gemini_portion_model,
        system=repair_system,
        prompt=repair_prompt,
        temperature=0.0,
        max_output_tokens=1400,
    )
    obj = model_json_or_400(repaired)
    return obj if isinstance(obj, dict) else {}


async def identify_repair_if_list(raw: Any, hints: Optional[List[str]]) -> Dict[str, Any]:
    if isinstance(raw, list):
        repair_out = await gemini.generate_text(
            model=settings.gemini_classifier_model,
            system=IDENTIFY_REPAIR_SYSTEM,
            prompt=identify_repair_prompt(raw, hints),
            temperature=0.2,
            max_output_tokens=700,
        )
        repaired = model_json_or_400(repair_out)
        if not isinstance(repaired, dict):
            raise HTTPException(400, "Identify repair did not return a JSON object")
        return repaired

    if isinstance(raw, dict):
        return raw

    raise HTTPException(400, "Invalid identify output type")

def normalize_identify_dict(obj: Dict[str, Any]) -> IdentifyResponse:
    candidates = obj.get("candidates") or []
    chosen = obj.get("chosen")

    if isinstance(chosen, str):
        chosen_lower = chosen.strip().lower()
        match = None
        for c in candidates:
            if not isinstance(c, dict):
                continue
            if chosen_lower in (
                str(c.get("name", "")).lower(),
                str(c.get("normalized_name", "")).lower()
            ):
                match = c
                break
        if match is None:
            match = {
                "name": chosen,
                "confidence": 0.5,
                "normalized_name": chosen,
                "cuisine": None,
                "is_packaged": None,
                "notes": "chosen string -> synthesized",
            }
            candidates.append(match)
        obj["chosen"] = match
        obj["candidates"] = candidates

    if isinstance(obj.get("chosen"), dict) and not obj.get("candidates"):
        obj["candidates"] = [obj["chosen"]]

    return IdentifyResponse.model_validate(obj)


def normalize_identify_obj(obj: Any) -> Dict[str, Any]:
    if not isinstance(obj, dict):
        return {"candidates": [], "chosen": {"name": "unknown", "confidence": 0.0, "normalized_name": "unknown"}}

    candidates = obj.get("candidates") or []
    chosen = obj.get("chosen")

    # handle chosen sometimes being string
    if isinstance(chosen, str):
        chosen_str = chosen
        chosen = None
        for c in candidates:
            if isinstance(c, dict) and (c.get("name") == chosen_str or c.get("normalized_name") == chosen_str):
                chosen = c
                break
        if chosen is None:
            chosen = {"name": chosen_str, "confidence": 0.5, "normalized_name": chosen_str}

    if not isinstance(chosen, dict):
        # fall back to best candidate
        best = None
        best_conf = -1
        for c in candidates:
            if isinstance(c, dict):
                conf = float(c.get("confidence") or 0.0)
                if conf > best_conf:
                    best_conf = conf
                    best = c
        chosen = best or {"name": "unknown", "confidence": 0.0, "normalized_name": "unknown"}

    # ensure candidates list of dicts
    clean_candidates = [c for c in candidates if isinstance(c, dict)]
    if not clean_candidates:
        clean_candidates = [chosen]

    obj["candidates"] = clean_candidates[:3]
    obj["chosen"] = chosen
    return obj

def normalize_portion_obj(obj: Any) -> Dict[str, Any]:
    if not isinstance(obj, dict):
        obj = {}

    obj.setdefault("servings", 1.0)
    obj.setdefault("grams_total", 0.0)
    obj.setdefault("items_count", None)
    obj.setdefault("household", None)
    obj.setdefault("confidence", 0.6)
    obj.setdefault("assumptions", [])

    # household must be string or None
    h = obj.get("household")
    if h is not None and not isinstance(h, str):
        obj["household"] = str(h)

    if not isinstance(obj.get("assumptions"), list):
        obj["assumptions"] = [str(obj["assumptions"])]

    try:
        obj["servings"] = float(obj["servings"] or 1.0)
    except Exception:
        obj["servings"] = 1.0

    try:
        obj["grams_total"] = float(obj["grams_total"] or 0.0)
    except Exception:
        obj["grams_total"] = 0.0

    ic = obj.get("items_count")
    if ic is None:
        obj["items_count"] = None
    else:
        try:
            obj["items_count"] = float(ic)
        except Exception:
            obj["items_count"] = None

    try:
        obj["confidence"] = float(obj["confidence"])
    except Exception:
        obj["confidence"] = 0.6

    # clamp
    obj["confidence"] = max(0.0, min(1.0, obj["confidence"]))
    return obj

def normalize_nutrients_obj(obj: Any, req: NutrientsRequest) -> Dict[str, Any]:
    if not isinstance(obj, dict):
        obj = {}

    # force required top-level keys
    obj.setdefault("food_name", req.food_name)
    if "portion" not in obj or not isinstance(obj["portion"], dict):
        # use request portion as source of truth
        p = req.portion
        obj["portion"] = {
            "servings": p.servings,
            "grams_total": p.grams_total,
            "items_count": p.items_count,
            "household": p.household,
            "confidence": getattr(p, "confidence", 0.6) or 0.6,
            "assumptions": getattr(p, "assumptions", []) or [],
        }

    obj.setdefault("calories_kcal", 0.0)
    obj.setdefault("macros", {})
    obj.setdefault("micros", {})
    obj.setdefault("vitamins", {})
    obj.setdefault("minerals", {})
    obj.setdefault("ingredients_guess", [])
    obj.setdefault("allergens_guess", [])
    obj.setdefault("data_sources", [])
    obj.setdefault("notes", [])

    # lists must be lists
    for k in ("ingredients_guess", "allergens_guess", "data_sources", "notes"):
        if not isinstance(obj.get(k), list):
            obj[k] = [str(obj.get(k))]

    # If user doesn't want per 100g
    if not req.include_per_100g:
        for grp in ("macros", "micros", "vitamins", "minerals"):
            if isinstance(obj.get(grp), dict):
                for _, item in obj[grp].items():
                    if isinstance(item, dict):
                        item["per_100g_amount"] = None

    return obj

# ---------------- Routes ----------------

@app.get("/health")
async def health():
    return {"ok": True, "env": settings.app_env}

@app.get("/v1/meta")
async def meta():
    return {
        "version": "1.0.0",
        "models": {
            "gemini_classifier": settings.gemini_classifier_model,
            "gemini_portion": settings.gemini_portion_model,
            "perplexity_sonar": settings.pplx_sonar_model,
        },
    }

@app.post("/v1/food/identify", response_model=IdentifyResponse)
async def identify_food(req: IdentifyRequest):
    if req.mode != "text" or not req.text:
        raise HTTPException(status_code=422, detail="Use mode=text with text, or use /v1/food/identify-image.")

    out = await gemini.generate_text(
        model=settings.gemini_classifier_model,
        system=IDENTIFY_SYSTEM,
        prompt=identify_user_prompt_text(req.text, req.hints),
        temperature=0.2,
        max_output_tokens=800,
    )

    raw = model_json_or_400(out)
    obj = normalize_identify_obj(raw)
    return IdentifyResponse.model_validate(obj)

@app.post("/v1/food/identify-image", response_model=IdentifyResponse)
async def identify_food_image(file: UploadFile = File(...), hints: Optional[str] = None):
    image_bytes = await file.read()
    if not image_bytes:
        raise HTTPException(status_code=422, detail="Empty upload (image_bytes is empty).")

    hints_list = [h.strip() for h in (hints or "").split(",") if h.strip()] or None

    out = await gemini.generate_with_image(
        model=settings.gemini_classifier_model,
        system=IDENTIFY_SYSTEM,
        prompt=identify_user_prompt_image(hints_list),
        image_bytes=image_bytes,
        mime_type=file.content_type or "image/jpeg",
        temperature=0.2,
        max_output_tokens=900,
    )

    raw = model_json_or_400(out)
    obj = normalize_identify_obj(raw)
    return IdentifyResponse.model_validate(obj)

@app.post("/v1/food/portion", response_model=PortionResponse)
async def portion_text(req: PortionRequest):
    if req.mode != "text":
        raise HTTPException(status_code=422, detail="Use /v1/food/portion-image for images.")

    out = await gemini.generate_text(
        model=settings.gemini_portion_model,
        system=PORTION_SYSTEM,
        prompt=portion_prompt_text(req.food_name, req.text_context),
        temperature=0.2,
        max_output_tokens=700,
    )

    try:
        raw = model_json_or_400(out)
    except HTTPException:
        raw = await force_json_with_gemini(
            '{"servings":number,"grams_total":number,"items_count":number|null,"household":string|null,"confidence":number,"assumptions":string[]}',
            out,
        )

    raw = normalize_portion_obj(raw)
    est = PortionEstimate.model_validate(raw)

    if est.grams_total <= 0:
        est = PortionEstimate(
            servings=max(1.0, est.servings),
            grams_total=100.0,
            items_count=None,
            household=req.household_measure or "1 serving (default 100g)",
            confidence=0.3,
            assumptions=["Defaulted to 100g because model output was uncertain"],
        )

    return PortionResponse(food_name=req.food_name, portion=est)

@app.post("/v1/food/portion-image", response_model=PortionResponse)
async def portion_image(food_name: str, file: UploadFile = File(...), text_context: Optional[str] = None):
    image_bytes = await file.read()
    if not image_bytes:
        raise HTTPException(status_code=422, detail="Empty upload (image_bytes is empty).")

    out = await gemini.generate_with_image(
        model=settings.gemini_portion_model,
        system=PORTION_SYSTEM,
        prompt=portion_prompt_image(food_name, text_context),
        image_bytes=image_bytes,
        mime_type=file.content_type or "image/jpeg",
        temperature=0.2,
        max_output_tokens=900,
    )

    try:
        raw = model_json_or_400(out)
    except HTTPException:
        raw = await force_json_with_gemini(
            '{"servings":number,"grams_total":number,"items_count":number|null,"household":string|null,"confidence":number,"assumptions":string[]}',
            out,
        )

    raw = normalize_portion_obj(raw)
    est = PortionEstimate.model_validate(raw)

    if est.grams_total <= 0:
        est = PortionEstimate(
            servings=max(1.0, est.servings),
            grams_total=100.0,
            items_count=None,
            household="1 serving (default 100g)",
            confidence=0.3,
            assumptions=["Defaulted to 100g because model output was uncertain"],
        )

    return PortionResponse(food_name=food_name, portion=est)

@app.post("/v1/food/nutrients", response_model=NutrientsResponse)
async def nutrients(req: NutrientsRequest):
    if req.portion is None:
        raise HTTPException(status_code=422, detail="portion is required. Call /v1/food/portion first.")
    return await compute_nutrients(req)

@app.post("/v1/food/analyze", response_model=AnalyzeResponse)
async def analyze_text(req: AnalyzeTextRequest):
    # 1) Identify
    identify_out = await gemini.generate_text(
        model=settings.gemini_classifier_model,
        system=IDENTIFY_SYSTEM,
        prompt=identify_user_prompt_text(req.text, req.hints),
        temperature=0.2,
        max_output_tokens=800,
    )
    identify_raw = model_json_or_400(identify_out)
    identify_obj = normalize_identify_obj(identify_raw)
    identify_res = IdentifyResponse.model_validate(identify_obj)
    chosen = identify_res.chosen.normalized_name or identify_res.chosen.name

    # 2) Portion (text)
    portion_out = await gemini.generate_text(
        model=settings.gemini_portion_model,
        system=PORTION_SYSTEM,
        prompt=portion_prompt_text(chosen, req.text),
        temperature=0.2,
        max_output_tokens=700,
    )
    try:
        portion_raw = model_json_or_400(portion_out)
    except HTTPException:
        portion_raw = await force_json_with_gemini(
            '{"servings":number,"grams_total":number,"items_count":number|null,"household":string|null,"confidence":number,"assumptions":string[]}',
            portion_out,
        )

    portion_raw = normalize_portion_obj(portion_raw)
    portion_est = PortionEstimate.model_validate(portion_raw)

    if portion_est.grams_total <= 0:
        portion_est = PortionEstimate(
            servings=1.0,
            grams_total=100.0,
            items_count=None,
            household="1 serving (default 100g)",
            confidence=0.3,
            assumptions=["Defaulted to 100g because model output was uncertain"],
        )

    portion_res = PortionResponse(food_name=chosen, portion=portion_est)

    # 3) Nutrients
    nreq = NutrientsRequest(
        food_name=chosen,
        portion=portion_est,
        region=req.region,
        include_per_100g=req.include_per_100g,
    )

    nutrients_out = await pplx.chat(
        model=settings.pplx_sonar_model,
        system=NUTRIENTS_SYSTEM,
        user=nutrients_prompt(nreq),
        temperature=0.2,
        max_tokens=900,
    )

    try:
        nutrients_raw = model_json_or_400(nutrients_out)
        if not isinstance(nutrients_raw, dict):
            raise HTTPException(status_code=400, detail="Nutrients output not an object")
    except HTTPException:
        nutrients_raw = await force_json_with_gemini("NutrientsResponse schema JSON object", nutrients_out)

    nutrients_raw = normalize_nutrients_obj(nutrients_raw, nreq)
    nutrients_res = NutrientsResponse.model_validate(nutrients_raw)

    return AnalyzeResponse(
        identify=identify_res,
        portion=portion_res,
        nutrients=nutrients_res,
        cost_tier={"identify": "$", "portion": "$$", "nutrients": "$$$$"},
    )
@app.post("/v1/food/analyze-image", response_model=AnalyzeResponse)
async def analyze_image(
    file: UploadFile = File(...),
    hints: Optional[str] = None,
    region: str = "IN",
    include_per_100g: bool = True,
):
    image_bytes = await file.read()
    if not image_bytes:
        raise HTTPException(400, "Uploaded image is empty")

    hints_list = [h.strip() for h in (hints or "").split(",") if h.strip()] or None

    identify_out = await gemini.generate_with_image(
        model=settings.gemini_classifier_model,
        system=IDENTIFY_SYSTEM,
        prompt=identify_user_prompt_image(hints_list),
        image_bytes=image_bytes,
        mime_type=file.content_type or "image/jpeg",
        temperature=0.2,
        max_output_tokens=900,
    )

    raw_ident = model_json_or_400(identify_out)
    raw_ident = await identify_repair_if_list(raw_ident, hints_list)  # keep your existing repair
    identify = normalize_identify_dict(raw_ident)                     # keep your existing normalize that returns IdentifyResponse
    food_name = identify.chosen.normalized_name or identify.chosen.name

    portion_out = await gemini.generate_with_image(
        model=settings.gemini_portion_model,
        system=PORTION_SYSTEM,
        prompt=portion_prompt_image(food_name, ctx=None),
        image_bytes=image_bytes,
        mime_type=file.content_type or "image/jpeg",
        temperature=0.2,
        max_output_tokens=900,
    )

    raw_portion = model_json_or_400(portion_out)
    if not isinstance(raw_portion, dict):
        raw_portion = await force_json_with_gemini(
            '{"servings":number,"grams_total":number,"items_count":number|null,"household":string|null,"confidence":number,"assumptions":string[]}',
            portion_out,
        )

    raw_portion = normalize_portion_obj(raw_portion)
    portion = PortionEstimate.model_validate(raw_portion)

    nreq = NutrientsRequest(
        food_name=food_name,
        portion=portion,
        region=region,
        include_per_100g=include_per_100g,
    )

    nutrients_res = await compute_nutrients(nreq)

    return AnalyzeResponse(
        identify=identify,
        portion=PortionResponse(food_name=food_name, portion=portion),
        nutrients=nutrients_res,
        cost_tier={"identify": "$", "portion": "$$", "nutrients": "$$$$"},
    )

# --------------------------------------------------
# PPT Routes (single instance)
# --------------------------------------------------

@app.post("/v1/ppt/generate", response_model=PPTGenerateResponse)
async def generate_ppt(req: PPTGenerateRequest):
    topic = (req.topic or "").strip()
    if not topic:
        raise HTTPException(422, "topic is required")

    design_number = req.design_number
    if design_number < 1 or design_number > 7:
        design_number = 1

    safe_topic = re.sub(r"[^\w\s\.\-\(\)]", "", topic).replace("\n", "").strip()
    filename = f"{safe_topic[:40]}_{uuid.uuid4().hex}"

    deck_text = await generate_presentation_text_async(safe_topic)
    fallback_used = False
    if not deck_text:
        deck_text = _local_fallback_presentation(safe_topic)
        fallback_used = True

    ppt_path = await create_presentation_async(deck_text, design_number, filename)

    base_url = (settings.public_base_url or "").rstrip("/")
    if not base_url:
        # Render provides your onrender URL; easiest is to set PUBLIC_BASE_URL env var in Render
        # Example: https://your-service.onrender.com
        base_url = "https://api.dashovia.com"

    download_url = f"{base_url}/v1/ppt/download/{Path(ppt_path).name}"

    return PPTGenerateResponse(
        status="success",
        filename=Path(ppt_path).name,
        download_url=download_url,
        fallback_used=fallback_used,
    )

@app.get("/v1/ppt/download/{filename}")
async def download_ppt(filename: str):
    file_path = GENERATED_DIR / filename
    if not file_path.is_file():
        raise HTTPException(404, "File not found")
    return FileResponse(
        file_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=filename,
    )
