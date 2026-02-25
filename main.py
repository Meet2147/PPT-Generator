

import os
import re
import uuid
import random
import asyncio
import json
import hashlib
import importlib.util
import sys
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List

import requests
import g4f
from g4f.client import Client
from pptx import Presentation
from pptx.util import Inches

from fastapi import Depends, FastAPI, HTTPException, Query, Request, Response
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
import uvicorn

# ---------------------------
# Prompt Template
# ---------------------------
PROMPT_TEMPLATE = """Write a presentation/powerpoint about the user's topic. You only answer with the presentation. Follow the structure of the example.

Notice:
- You do all the presentation text for the user.
- You write the texts no longer than 250 characters!
- You make very short titles!
- You make the presentation easy to understand.
- The presentation has a table of contents.
- The presentation has a summary.
- At least 7 slides.
- For each slide, after the #Content: line, add an #Image: line describing a relevant image that could visually represent the slide's topic.
- If no image is relevant, write #Image: none.

Example! - Stick to this formatting exactly!
#Title: TITLE OF THE PRESENTATION

#Slide: 1
#Header: table of contents
#Content: 1. CONTENT OF THIS POWERPOINT
2. CONTENTS OF THIS POWERPOINT
3. CONTENT OF THIS POWERPOINT
#Image: a 3D illustration of a table of contents in a book

#Slide: 2
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE
#Image: relevant illustration description here

#Slide: 3
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE
#Image: relevant illustration description here

#Slide: 4
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE
#Image: relevant illustration description here

#Slide: 5
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE
#Image: relevant illustration description here

#Slide: 6
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE
#Image: relevant illustration description here

#Slide: 7
#Header: summary
#Content: CONTENT OF THE SUMMARY
#Image: an illustration of a person reviewing a summary report

#Slide: END"""

# ---------------------------
# Clients + App Setup
# ---------------------------
# Keep g4f only for images (you asked to move text gen to Sonar)
g4f_client = Client()

app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],   # tighten for prod
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

Path("GeneratedPresentations").mkdir(exist_ok=True)
Path("Cache").mkdir(exist_ok=True)
Path("Designs").mkdir(exist_ok=True)

# ---------------------------
# Perplexity (Sonar) minimal client
# ---------------------------
PPLX_API_KEY = os.getenv("PPLX_API_KEY")
PPLX_URL = "https://api.perplexity.ai/chat/completions"
PPLX_MODEL = "sonar-pro"  # tokens-only (no *-online); switch to "sonar" if you prefer

def _perplexity_chat(messages, model=PPLX_MODEL, max_tokens=1200, temperature=0.4):
    if not PPLX_API_KEY:
        raise RuntimeError(
            "Missing PPLX_API_KEY. Set it in your environment before running the server."
        )
    payload = {
        "model": model,
        "messages": messages,
        "temperature": temperature,
        "top_p": 0.95,
        "max_tokens": max_tokens,
        "presence_penalty": 0.0,
        "frequency_penalty": 0.0,
        # ensure token-only text generation (no images or search)
        "return_images": False,
        "stream": False,
    }
    headers = {
        "Authorization": f"Bearer {PPLX_API_KEY}",
        "Content-Type": "application/json",
    }
    resp = requests.post(PPLX_URL, json=payload, headers=headers, timeout=60)
    if resp.status_code != 200:
        raise RuntimeError(f"Perplexity API error {resp.status_code}: {resp.text}")
    data = resp.json()
    # Perplexity follows OpenAI-like schema
    content = data.get("choices", [{}])[0].get("message", {}).get("content", "")
    if not isinstance(content, str):
        content = str(content or "")
    return content

# ---------------------------
# Helpers
# ---------------------------
def _local_fallback_presentation(topic: str) -> str:
    topic_title = (topic or "Your Topic").title()
    return f"""#Title: {topic_title}

#Slide: 1
#Header: table of contents
#Content: 1. Intro
2. Why it matters
3. Key points
4. Examples
5. Tips
6. Pitfalls
7. Summary
#Image: a 3D illustration of a table of contents in a book

#Slide: 2
#Header: intro
#Content: Brief overview of {topic}. Scope and goal of this talk.
#Image: relevant illustration of the topic

#Slide: 3
#Header: why it matters
#Content: Impact, use-cases, and benefits in simple terms.
#Image: simple infographic with benefits

#Slide: 4
#Header: key points
#Content: 3–5 core ideas about {topic}, kept concise.
#Image: icons representing each key idea

#Slide: 5
#Header: examples
#Content: 2–3 quick examples showing the concept in action.
#Image: storyboard-like illustration

#Slide: 6
#Header: tips & pitfalls
#Content: Do's and don'ts for better results.
#Image: checklist illustration

#Slide: 7
#Header: summary
#Content: Short recap and practical next steps.
#Image: an illustration of a person reviewing a summary report

#Slide: END
"""

async def generate_image_async(prompt: str) -> str | None:
    if not prompt or prompt.strip().lower() == "none":
        return None
    try:
        resp = await asyncio.to_thread(
            g4f_client.images.generate,
            model="flux",             # keep if working for you; otherwise stub this out
            prompt=prompt,
            response_format="url"
        )
        if hasattr(resp, "data") and resp.data:
            return resp.data[0].url
    except Exception as e:
        print(f"[generate_image_async] Error generating image: {e}")
    return None

async def generate_presentation_text_async(user_input: str) -> str:
    """Use Perplexity Sonar (no web/online) to generate strict deck content."""
    system_msg = PROMPT_TEMPLATE
    user_msg = f"The user wants a presentation about {user_input}"

    # Attempt 1
    try:
        content = await asyncio.to_thread(
            _perplexity_chat,
            [
                {"role": "system", "content": system_msg},
                {"role": "user", "content": user_msg},
            ],
            PPLX_MODEL,
            1400,   # allow enough tokens for at least 7 slides
            0.35,
        )
        if "#Slide:" not in content or "#Title:" not in content:
            raise ValueError("Model response missing required markers (#Title:/#Slide:).")
        return content
    except Exception as e:
        print(f"[generate_presentation_text_async] Sonar first attempt failed: {e}")

    # Attempt 2 with stricter instruction
    try:
        strict_prompt = "ONLY output in the exact template below. Do NOT add extra text.\n\n" + PROMPT_TEMPLATE
        content = await asyncio.to_thread(
            _perplexity_chat,
            [
                {"role": "system", "content": strict_prompt},
                {"role": "user", "content": user_msg},
            ],
            PPLX_MODEL,
            1400,
            0.2,
        )
        if "#Slide:" not in content or "#Title:" not in content:
            raise ValueError("Model response missing required markers on retry.")
        return content
    except Exception as e:
        print(f"[generate_presentation_text_async] Sonar retry failed: {e}")
        return ""

async def create_presentation_async(text_content: str, design_number: int, presentation_name: str):
    # Open template (fallback to Design-1, else blank)
    template_path = Path(f"Designs/Design-{design_number}.pptx")
    if not template_path.exists():
        template_path = Path("Designs/Design-1.pptx")

    if template_path.exists():
        try:
            presentation = Presentation(str(template_path))
        except Exception as e:
            print(f"[create_presentation_async] Could not open template, using blank. Error: {e}")
            presentation = Presentation()
    else:
        presentation = Presentation()

    # State
    slide_count = 0
    slide_title = ""
    slide_content = ""
    slide_image_prompt = None
    last_slide_layout_index = -1
    first_time = True

    lines = [ln.rstrip("\n") for ln in text_content.splitlines()]
    i = 0

    # Layout helpers
    def _safe_layout(idx: int) -> int:
        if len(presentation.slide_layouts) == 0:
            return 0
        return max(0, min(idx, len(presentation.slide_layouts) - 1))

    slide_layout_index = _safe_layout(1 if len(presentation.slide_layouts) > 1 else 0)
    slide_placeholder_index = 2 if slide_layout_index == 8 else (1 if len(presentation.slide_layouts) > 1 else 0)

    # Commit helper
    async def commit_slide():
        nonlocal slide_title, slide_content, slide_image_prompt, slide_count
        if slide_count > 0 and (slide_title or slide_content):
            slide = presentation.slides.add_slide(presentation.slide_layouts[_safe_layout(slide_layout_index)])
            # Title
            try:
                slide.shapes.title.text = slide_title
            except Exception:
                try:
                    slide.placeholders[0].text = slide_title
                except Exception:
                    pass
            # Body
            try:
                body_shape = slide.shapes.placeholders[slide_placeholder_index]
                if hasattr(body_shape, "text_frame"):
                    body_shape.text_frame.text = slide_content
                else:
                    body_shape.text = slide_content
            except Exception:
                pass

            # Image (if any)
            if slide_image_prompt:
                img_url = await generate_image_async(slide_image_prompt)
                if img_url:
                    try:
                        resp = await asyncio.to_thread(requests.get, img_url, timeout=20)
                        resp.raise_for_status()
                        image_stream = BytesIO(resp.content)
                        slide.shapes.add_picture(image_stream, Inches(5), Inches(1.5), width=Inches(4))
                    except Exception as e:
                        print(f"[create_presentation_async] Could not add image: {e}")

    # Parse
    while i < len(lines):
        line = lines[i]

        if line.startswith("#Title:"):
            slide_title_text = line.replace("#Title:", "").strip()
            title_layout = _safe_layout(0)
            slide = presentation.slides.add_slide(presentation.slide_layouts[title_layout])
            try:
                slide.shapes.title.text = slide_title_text
            except Exception:
                try:
                    slide.placeholders[0].text = slide_title_text
                except Exception:
                    pass
            i += 1
            continue

        if line.startswith("#Slide:"):
            await commit_slide()
            slide_content = ""
            slide_image_prompt = None
            slide_title = ""
            slide_count += 1

            if len(presentation.slide_layouts) >= 9:
                layout_choices = [1, 7, 8]
            else:
                layout_choices = [1] if len(presentation.slide_layouts) > 1 else [0]

            if first_time:
                slide_layout_index = _safe_layout(1 if len(presentation.slide_layouts) > 1 else 0)
                slide_placeholder_index = 2 if slide_layout_index == 8 else (1 if len(presentation.slide_layouts) > 1 else 0)
                first_time = False
            else:
                next_idx = last_slide_layout_index
                tries = 0
                while next_idx == last_slide_layout_index and tries < 10:
                    next_idx = random.choice(layout_choices)
                    tries += 1
                slide_layout_index = _safe_layout(next_idx)
                slide_placeholder_index = 2 if slide_layout_index == 8 else (1 if len(presentation.slide_layouts) > 1 else 0)

            last_slide_layout_index = slide_layout_index
            i += 1
            continue

        if line.startswith("#Header:"):
            slide_title = line.replace("#Header:", "").strip()
            i += 1
            continue

        if line.startswith("#Content:"):
            slide_content = line.replace("#Content:", "").strip()
            i += 1
            while i < len(lines) and not lines[i].startswith("#"):
                slide_content += "\n" + lines[i]
                i += 1
            continue

        if line.startswith("#Image:"):
            slide_image_prompt = line.replace("#Image:", "").strip()
            i += 1
            continue

        i += 1

    # Final commit and save
    await commit_slide()

    out_dir = Path("GeneratedPresentations")
    out_dir.mkdir(parents=True, exist_ok=True)
    file_path = out_dir / f"{presentation_name}.pptx"
    await asyncio.to_thread(presentation.save, str(file_path))
    return str(file_path)

# ---------------------------
# Routes
# ---------------------------
@app.get("/auto-download/{filename}")
async def download_file(filename: str):
    file_path = Path("GeneratedPresentations") / filename
    if file_path.is_file():
        return FileResponse(
            file_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=filename
        )
    else:
        return Response(content="File not found", status_code=404)

@app.post("/generate-ppt/")
async def generate_ppt(user_input: str):
    """
    POST /generate-ppt/?user_input=Your topic here
    """
    try:
        user_text = user_input or ""
        last_char = user_text[-1] if user_text else ""
        input_string = re.sub(r"[^\w\s\.\-\(\)]", "", user_text).replace("\n", "")
        design_number = 2

        # "topic 3" => design 3
        if last_char.isdigit():
            try:
                design_number = int(last_char)
                input_string = (
                    user_text[:-2].strip()
                    if len(user_text) >= 2 and user_text[-2].isspace()
                    else user_text[:-1].strip()
                )
                input_string = re.sub(r"[^\w\s\.\-\(\)]", "", input_string).replace("\n", "")
            except Exception:
                design_number = 2

        if design_number > 7 or design_number == 0:
            design_number = 1

        filename = f"{input_string}_{uuid.uuid4().hex}"

        # Generate with Sonar → strict format
        presentation_text = await generate_presentation_text_async(input_string)

        used_fallback = False
        if not presentation_text:
            print("[/generate-ppt] Sonar did not return valid content; using local fallback.")
            presentation_text = _local_fallback_presentation(input_string)
            used_fallback = True

        file_path = await create_presentation_async(presentation_text, design_number, filename)

        # TODO: change to your deployment host if different
        auto_download_url = f"https://ppt-generator-mtu0.onrender.com/auto-download/{Path(file_path).name}"

        return {
            "status": "success",
            "url": auto_download_url,
            "fallback_used": used_fallback
        }

    except Exception as e:
        return {"status": "error", "message": f"Server exception: {e}"}


# ---------------------------
# Workouts API (from external/workouts-premium)
# ---------------------------
WORKOUTS_ROOT = Path(__file__).parent / "external" / "workouts-premium"
USAGE_API_BASE_URL = os.getenv("USAGE_API_BASE_URL", "https://getmyworkouts.dashovia.com/api/v1").rstrip("/")
USAGE_API_TIMEOUT_SECONDS = float(os.getenv("USAGE_API_TIMEOUT_SECONDS", "8"))
WORKOUT_GIF_DIR = WORKOUTS_ROOT / "gifs_360x360"
WORKOUT_GIF_PROXY_ROUTE = "/workouts/media/gif"


def _load_workouts_json(filename: str) -> list[dict]:
    path = WORKOUTS_ROOT / filename
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


WORKOUTS_AVAILABLE = WORKOUTS_ROOT.exists()
WORKOUTS_IMPORT_ERROR = ""
WORKOUT_EXERCISES: list[dict] = []
WORKOUT_MUSCLES: list[dict] = []
WORKOUT_BODY_PARTS: list[dict] = []
WORKOUT_MUSCLE_NAMES: dict[str, str] = {}
WORKOUT_BODY_PART_NAMES: dict[str, str] = {}
WORKOUT_TARGET_MUSCLE_SET: set[str] = set()
workout_analyze_fn = None

if WORKOUTS_AVAILABLE:
    try:
        WORKOUT_EXERCISES = _load_workouts_json("exercises.json")
        WORKOUT_MUSCLES = _load_workouts_json("muscles.json")
        WORKOUT_BODY_PARTS = _load_workouts_json("bodyParts.json")
        WORKOUT_MUSCLE_NAMES = {m["name"].lower(): m["name"] for m in WORKOUT_MUSCLES}
        WORKOUT_BODY_PART_NAMES = {b["name"].lower(): b["name"] for b in WORKOUT_BODY_PARTS}
        WORKOUT_TARGET_MUSCLE_SET = {
            tm.lower() for e in WORKOUT_EXERCISES for tm in e.get("targetMuscles", [])
        }

        spec = importlib.util.spec_from_file_location(
            "workouts_perplexity_service", WORKOUTS_ROOT / "perplexity_service.py"
        )
        if spec and spec.loader:
            module = importlib.util.module_from_spec(spec)
            spec.loader.exec_module(module)
            workout_analyze_fn = getattr(module, "analyze_workout_question", None)
    except Exception as exc:
        WORKOUTS_AVAILABLE = False
        WORKOUTS_IMPORT_ERROR = str(exc)


WORKOUT_TARGET_MUSCLE_ALIASES = {
    "deltoid": "delts",
    "deltoids": "delts",
    "shoulder": "delts",
    "shoulders": "delts",
    "chest": "pectorals",
    "pec": "pectorals",
    "pecs": "pectorals",
    "back": "upper back",
}


def _normalize_target_muscle(value: str) -> str:
    normalized = value.lower().strip()
    if normalized in WORKOUT_TARGET_MUSCLE_SET:
        return normalized
    return WORKOUT_TARGET_MUSCLE_ALIASES.get(normalized, normalized)


def _workout_gif_url(gif_filename: str) -> str:
    return f"{WORKOUT_GIF_PROXY_ROUTE}/{gif_filename}"


def _resolve_api_key(request: Request) -> str:
    header_key = request.headers.get("x-api-key", "").strip()
    if header_key:
        return header_key
    auth = request.headers.get("authorization", "").strip()
    if auth.lower().startswith("bearer "):
        return auth.split(" ", 1)[1].strip()
    return ""


def _usage_api_error_detail(response: requests.Response) -> str:
    try:
        payload = response.json()
        if isinstance(payload, dict):
            return str(payload.get("detail") or payload.get("message") or payload)
        return str(payload)
    except Exception:
        return response.text or "Usage service request failed."


def _validate_api_key_with_usage_service(api_key: str) -> None:
    try:
        response = requests.get(
            f"{USAGE_API_BASE_URL}/usage",
            headers={"x-api-key": api_key, "accept": "application/json"},
            timeout=USAGE_API_TIMEOUT_SECONDS,
        )
    except requests.RequestException:
        raise HTTPException(
            status_code=503,
            detail="Unable to reach usage service for API key validation.",
        )

    if response.status_code == 401:
        raise HTTPException(status_code=401, detail="Missing or invalid API key.")
    if response.status_code >= 400:
        raise HTTPException(
            status_code=502,
            detail=f"Usage service validation failed: {_usage_api_error_detail(response)}",
        )


def _consume_usage_with_service(api_key: str, usage_type: str) -> None:
    endpoint = "consume-ai" if usage_type == "ai" else "consume-standard"
    try:
        response = requests.post(
            f"{USAGE_API_BASE_URL}/{endpoint}",
            headers={"x-api-key": api_key, "accept": "application/json"},
            timeout=USAGE_API_TIMEOUT_SECONDS,
        )
    except requests.RequestException:
        raise HTTPException(status_code=503, detail="Unable to reach usage service.")

    if response.status_code == 401:
        raise HTTPException(status_code=401, detail="Missing or invalid API key.")
    if response.status_code == 429:
        raise HTTPException(status_code=429, detail=_usage_api_error_detail(response))
    if response.status_code >= 400:
        raise HTTPException(
            status_code=502,
            detail=f"Usage service consume failed: {_usage_api_error_detail(response)}",
        )


def _require_api_key(request: Request) -> str:
    api_key = _resolve_api_key(request)
    if not api_key:
        raise HTTPException(status_code=401, detail="Missing API key.")
    _validate_api_key_with_usage_service(api_key)
    return api_key


@app.get(f"{WORKOUT_GIF_PROXY_ROUTE}/" + "{gif_filename}")
async def workouts_get_gif(gif_filename: str, api_key: str = Depends(_require_api_key)):
    if not WORKOUTS_AVAILABLE:
        raise HTTPException(status_code=503, detail=f"Workouts API unavailable: {WORKOUTS_IMPORT_ERROR}")
    safe_name = Path(gif_filename).name
    if safe_name != gif_filename or not safe_name.lower().endswith(".gif"):
        raise HTTPException(status_code=400, detail="Invalid GIF filename")

    gif_path = WORKOUT_GIF_DIR / safe_name
    if not gif_path.exists():
        raise HTTPException(status_code=404, detail="GIF not found")

    _consume_usage_with_service(api_key, "standard")
    return FileResponse(path=gif_path, media_type="image/gif")


@app.get("/workouts")
async def workouts_root():
    return {
        "message": "Welcome to Exercise Database API",
        "endpoints": {
            "/workouts/available-muscles": "Get all available muscles",
            "/workouts/available-body-parts": "Get all available body parts",
            "/workouts/exercises-by-muscles": "Get exercises by muscle groups (select 2)",
            "/workouts/exercises-by-body-parts": "Get exercises by body parts (select 2)",
        },
    }


@app.get("/workouts/available-muscles", response_model=List[str])
async def workouts_available_muscles(api_key: str = Depends(_require_api_key)):
    if not WORKOUTS_AVAILABLE:
        raise HTTPException(status_code=503, detail=f"Workouts API unavailable: {WORKOUTS_IMPORT_ERROR}")
    result = sorted([m["name"] for m in WORKOUT_MUSCLES])
    _consume_usage_with_service(api_key, "standard")
    return result


@app.get("/workouts/available-body-parts", response_model=List[str])
async def workouts_available_body_parts(api_key: str = Depends(_require_api_key)):
    if not WORKOUTS_AVAILABLE:
        raise HTTPException(status_code=503, detail=f"Workouts API unavailable: {WORKOUTS_IMPORT_ERROR}")
    result = sorted([b["name"] for b in WORKOUT_BODY_PARTS])
    _consume_usage_with_service(api_key, "standard")
    return result


@app.get("/workouts/exercises-by-muscles")
async def workouts_exercises_by_muscles(
    muscle1: str = Query(..., description="First muscle group"),
    muscle2: str = Query(..., description="Second muscle group"),
    api_key: str = Depends(_require_api_key),
):
    if not WORKOUTS_AVAILABLE:
        raise HTTPException(status_code=503, detail=f"Workouts API unavailable: {WORKOUTS_IMPORT_ERROR}")
    m1 = muscle1.lower().strip()
    m2 = muscle2.lower().strip()
    if m1 not in WORKOUT_MUSCLE_NAMES or m2 not in WORKOUT_MUSCLE_NAMES:
        raise HTTPException(
            status_code=400,
            detail=f"Invalid muscle. Available muscles: {', '.join(sorted(WORKOUT_MUSCLE_NAMES.keys()))}",
        )

    m1_original = WORKOUT_MUSCLE_NAMES[m1]
    m2_original = WORKOUT_MUSCLE_NAMES[m2]
    ex1 = [
        {
            "exerciseId": e["exerciseId"],
            "name": e["name"],
            "gifUrl": _workout_gif_url(e["gifUrl"]),
            "equipment": e.get("equipments", []),
            "bodyPart": e.get("bodyParts", [None])[0],
            "secondaryMuscles": e.get("secondaryMuscles", []),
        }
        for e in WORKOUT_EXERCISES
        if m1_original in e.get("targetMuscles", [])
    ][:10]
    ex2 = [
        {
            "exerciseId": e["exerciseId"],
            "name": e["name"],
            "gifUrl": _workout_gif_url(e["gifUrl"]),
            "equipment": e.get("equipments", []),
            "bodyPart": e.get("bodyParts", [None])[0],
            "secondaryMuscles": e.get("secondaryMuscles", []),
        }
        for e in WORKOUT_EXERCISES
        if m2_original in e.get("targetMuscles", [])
    ][:10]

    _consume_usage_with_service(api_key, "standard")
    return {
        "selectedMuscles": [m1_original, m2_original],
        "results": {
            m1_original: {"count": len(ex1), "exercises": ex1},
            m2_original: {"count": len(ex2), "exercises": ex2},
        },
    }


@app.get("/workouts/exercises-by-body-parts")
async def workouts_exercises_by_body_parts(
    part1: str = Query(..., description="First body part"),
    part2: str = Query(..., description="Second body part"),
    api_key: str = Depends(_require_api_key),
):
    if not WORKOUTS_AVAILABLE:
        raise HTTPException(status_code=503, detail=f"Workouts API unavailable: {WORKOUTS_IMPORT_ERROR}")
    p1 = part1.lower().strip()
    p2 = part2.lower().strip()
    if p1 not in WORKOUT_BODY_PART_NAMES or p2 not in WORKOUT_BODY_PART_NAMES:
        raise HTTPException(
            status_code=400,
            detail=f"Invalid body part. Available body parts: {', '.join(sorted(WORKOUT_BODY_PART_NAMES.keys()))}",
        )

    p1_original = WORKOUT_BODY_PART_NAMES[p1]
    p2_original = WORKOUT_BODY_PART_NAMES[p2]
    ex1 = [
        {
            "exerciseId": e["exerciseId"],
            "name": e["name"],
            "gifUrl": _workout_gif_url(e["gifUrl"]),
            "equipment": e.get("equipments", []),
            "targetMuscles": e.get("targetMuscles", []),
            "secondaryMuscles": e.get("secondaryMuscles", []),
        }
        for e in WORKOUT_EXERCISES
        if p1_original in e.get("bodyParts", [])
    ][:10]
    ex2 = [
        {
            "exerciseId": e["exerciseId"],
            "name": e["name"],
            "gifUrl": _workout_gif_url(e["gifUrl"]),
            "equipment": e.get("equipments", []),
            "targetMuscles": e.get("targetMuscles", []),
            "secondaryMuscles": e.get("secondaryMuscles", []),
        }
        for e in WORKOUT_EXERCISES
        if p2_original in e.get("bodyParts", [])
    ][:10]

    _consume_usage_with_service(api_key, "standard")
    return {
        "selectedBodyParts": [p1_original, p2_original],
        "results": {
            p1_original: {"count": len(ex1), "exercises": ex1},
            p2_original: {"count": len(ex2), "exercises": ex2},
        },
    }


@app.get("/workouts/exercise/{exercise_id}")
async def workouts_exercise_details(exercise_id: str, api_key: str = Depends(_require_api_key)):
    if not WORKOUTS_AVAILABLE:
        raise HTTPException(status_code=503, detail=f"Workouts API unavailable: {WORKOUTS_IMPORT_ERROR}")
    exercise = next((e for e in WORKOUT_EXERCISES if e["exerciseId"] == exercise_id), None)
    if not exercise:
        raise HTTPException(status_code=404, detail="Exercise not found")

    _consume_usage_with_service(api_key, "standard")
    return {
        "exerciseId": exercise["exerciseId"],
        "name": exercise["name"],
        "gifUrl": _workout_gif_url(exercise["gifUrl"]),
        "targetMuscles": exercise.get("targetMuscles", []),
        "bodyParts": exercise.get("bodyParts", []),
        "equipments": exercise.get("equipments", []),
        "secondaryMuscles": exercise.get("secondaryMuscles", []),
        "instructions": exercise.get("instructions", []),
    }


@app.post("/workouts/suggest-workouts")
async def workouts_suggest_workouts(
    question: str = Query(..., description="User's workout question"),
    api_key: str = Depends(_require_api_key),
):
    if not WORKOUTS_AVAILABLE:
        raise HTTPException(status_code=503, detail=f"Workouts API unavailable: {WORKOUTS_IMPORT_ERROR}")
    if not workout_analyze_fn:
        raise HTTPException(status_code=503, detail="Workouts AI analyzer is unavailable.")

    try:
        muscle_list = [m["name"] for m in WORKOUT_MUSCLES]
        body_part_list = [b["name"] for b in WORKOUT_BODY_PARTS]
        analysis = workout_analyze_fn(question, muscle_list, body_part_list)

        exercises_list = []
        matched_any = False

        for muscle in analysis.get("suggested_muscles", []):
            normalized_muscle = _normalize_target_muscle(muscle)
            by_muscle = [
                {
                    "exerciseId": e["exerciseId"],
                    "name": e["name"],
                    "gifUrl": _workout_gif_url(e["gifUrl"]),
                    "equipment": e.get("equipments", []),
                    "targetMuscles": e.get("targetMuscles", []),
                    "secondaryMuscles": e.get("secondaryMuscles", []),
                    "bodyParts": e.get("bodyParts", []),
                }
                for e in WORKOUT_EXERCISES
                if normalized_muscle in [tm.lower() for tm in e.get("targetMuscles", [])]
            ][:8]
            if by_muscle:
                matched_any = True
            exercises_list.extend(by_muscle)

        if not matched_any:
            for body_part in analysis.get("suggested_body_parts", []):
                normalized_part = body_part.lower().strip()
                by_part = [
                    {
                        "exerciseId": e["exerciseId"],
                        "name": e["name"],
                        "gifUrl": _workout_gif_url(e["gifUrl"]),
                        "equipment": e.get("equipments", []),
                        "targetMuscles": e.get("targetMuscles", []),
                        "secondaryMuscles": e.get("secondaryMuscles", []),
                        "bodyParts": e.get("bodyParts", []),
                    }
                    for e in WORKOUT_EXERCISES
                    if normalized_part in [bp.lower() for bp in e.get("bodyParts", [])]
                ][:8]
                exercises_list.extend(by_part)

        seen = set()
        unique_exercises = []
        for ex in exercises_list:
            if ex["exerciseId"] in seen:
                continue
            seen.add(ex["exerciseId"])
            unique_exercises.append(ex)

        _consume_usage_with_service(api_key, "ai")
        return {
            "userQuestion": question,
            "analysis": analysis,
            "suggestedWorkouts": unique_exercises,
            "totalSuggestions": len(unique_exercises),
        }
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Error processing question: {exc}")


# ---------------------------
# Certiphy API (from external/Certiphy)
# ---------------------------
CERTIPHY_ROOT = Path(__file__).parent / "external" / "Certiphy"
CERTIPHY_AVAILABLE = CERTIPHY_ROOT.exists()
CERTIPHY_IMPORT_ERROR = ""
CERTIPHY_EXAMS = {}
certiphy_list_exams = None
certiphy_build_exam_like_quiz = None
certiphy_build_flashcards = None
certiphy_init_cache = None
certiphy_get_cached = None
certiphy_set_cached = None
CERTIPHY_SOURCES = {}

if CERTIPHY_AVAILABLE:
    try:
        sys.path.insert(0, str(CERTIPHY_ROOT))
        from app.registry import EXAMS as CERTIPHY_EXAMS, list_exams as certiphy_list_exams
        from app.services.quiz_builder import build_exam_like_quiz as certiphy_build_exam_like_quiz
        from app.services.flashcards_builder import (
            build_flashcards_from_quiz_questions as certiphy_build_flashcards,
        )
        from app.services.cache import (
            init_cache as certiphy_init_cache,
            get_cached as certiphy_get_cached,
            set_cached as certiphy_set_cached,
        )
        from app.services.sources.microsoft_learn import MicrosoftLearnSource
        from app.services.sources.aws_exam_guides_html import AwsExamGuidesHtmlSource
        from app.services.sources.google_cloud_html import GoogleCloudHtmlSource

        CERTIPHY_SOURCES = {
            "microsoft_learn": MicrosoftLearnSource(),
            "aws_exam_guides_html": AwsExamGuidesHtmlSource(),
            "google_cloud_html": GoogleCloudHtmlSource(),
        }
    except Exception as exc:
        CERTIPHY_AVAILABLE = False
        CERTIPHY_IMPORT_ERROR = str(exc)


CERTIPHY_DISCLAIMER = (
    "Practice questions are generated using AI APIs and may be inaccurate. "
    "They are for study assistance only, not a guarantee of exam results. "
    "Questions are original and not sourced from exam dumps."
)
CERTIPHY_CACHE_TTL_SECONDS = int(os.getenv("QUIZ_CACHE_TTL_SECONDS", str(30 * 24 * 3600)))
CERTIPHY_CACHE_ENABLED = os.getenv("QUIZ_CACHE_ENABLED", "true").lower() == "true"


def _certiphy_cache_key(exam_id: str, req: Dict[str, Any]) -> str:
    raw = {
        "exam_id": exam_id,
        "mode": req.get("mode", "practice"),
        "difficulty": req.get("difficulty", "easy"),
        "total_questions": int(req.get("total_questions", 25)),
        "scenario_ratio": float(req.get("scenario_ratio", 0.7)),
        "domain_distribution": req.get("domain_distribution"),
        "time_limit_minutes": req.get("time_limit_minutes") if req.get("mode", "practice") == "exam" else None,
    }
    return hashlib.sha1(json.dumps(raw, sort_keys=True).encode("utf-8")).hexdigest()


@app.on_event("startup")
def init_certiphy_cache():
    if CERTIPHY_AVAILABLE and CERTIPHY_CACHE_ENABLED and certiphy_init_cache:
        certiphy_init_cache()


@app.get("/certiphy/health")
def certiphy_health():
    if not CERTIPHY_AVAILABLE:
        return {"ok": False, "error": CERTIPHY_IMPORT_ERROR}
    return {"ok": True}


@app.get("/certiphy/exams")
def certiphy_exams():
    if not CERTIPHY_AVAILABLE or not certiphy_list_exams:
        raise HTTPException(status_code=503, detail=f"Certiphy API unavailable: {CERTIPHY_IMPORT_ERROR}")
    return {"exams": certiphy_list_exams()}


@app.get("/certiphy/exams/{exam_id}/blueprint")
async def certiphy_blueprint(exam_id: str):
    if not CERTIPHY_AVAILABLE:
        raise HTTPException(status_code=503, detail=f"Certiphy API unavailable: {CERTIPHY_IMPORT_ERROR}")
    exam_id = exam_id.upper().strip()
    cfg = CERTIPHY_EXAMS.get(exam_id)
    if not cfg:
        raise HTTPException(status_code=404, detail="Unknown exam_id. Use GET /certiphy/exams")

    src = CERTIPHY_SOURCES.get(cfg["source_type"])
    if not src:
        raise HTTPException(status_code=500, detail=f"No source handler for {cfg['source_type']}")

    sections = await src.fetch_sections(cfg["blueprint_url"])
    return {
        "exam_id": exam_id,
        "name": cfg["name"],
        "vendor": cfg["vendor"],
        "source": cfg["blueprint_url"],
        "sections": [{"index": i, "title": s["title"]} for i, s in enumerate(sections)],
    }


@app.post("/certiphy/exams/{exam_id}/generate-quiz")
async def certiphy_generate_quiz(exam_id: str, req: Dict[str, Any]):
    if not CERTIPHY_AVAILABLE or not certiphy_build_exam_like_quiz:
        raise HTTPException(status_code=503, detail=f"Certiphy API unavailable: {CERTIPHY_IMPORT_ERROR}")

    exam_id = exam_id.upper().strip()
    cfg = CERTIPHY_EXAMS.get(exam_id)
    if not cfg:
        raise HTTPException(status_code=404, detail="Unknown exam_id. Use GET /certiphy/exams")

    src = CERTIPHY_SOURCES.get(cfg["source_type"])
    if not src:
        raise HTTPException(status_code=500, detail=f"No source handler for {cfg['source_type']}")

    mode = req.get("mode", "practice")
    total_questions = int(req.get("total_questions", 25))
    difficulty = str(req.get("difficulty", cfg.get("difficulty_default", "easy"))).strip().lower()
    include_explanations = bool(req.get("include_explanations", True))
    scenario_ratio = float(req.get("scenario_ratio", 0.7))
    min_questions = int(req.get("min_questions", 10))
    max_retries = int(req.get("max_retries", 2))
    time_limit_minutes = req.get("time_limit_minutes")
    domain_distribution = req.get("domain_distribution")

    cache_key = _certiphy_cache_key(exam_id, req)
    if CERTIPHY_CACHE_ENABLED and certiphy_get_cached:
        cached = certiphy_get_cached(cache_key, ttl_seconds=CERTIPHY_CACHE_TTL_SECONDS)
        if cached:
            if not include_explanations and "questions" in cached:
                for q in cached["questions"]:
                    q["explanation"] = None
            return cached

    sections = await src.fetch_sections(cfg["blueprint_url"])
    questions = await certiphy_build_exam_like_quiz(
        exam_name=cfg["name"],
        vendor=cfg["vendor"],
        sections=sections,
        domains=cfg.get("domains") or {},
        keywords=cfg.get("keywords") or [],
        mode=mode,
        total_questions=total_questions,
        difficulty=difficulty,
        include_explanations=include_explanations,
        scenario_ratio=scenario_ratio,
        min_questions=min_questions,
        max_retries=max_retries,
        domain_distribution_override=domain_distribution,
    )

    if len(questions) < min_questions:
        raise HTTPException(
            status_code=502,
            detail=(
                f"Generated only {len(questions)} relevant questions for {exam_id}. "
                "Try increasing max_retries or total_questions."
            ),
        )

    response_obj = {
        "exam_id": exam_id,
        "exam_name": cfg["name"],
        "vendor": cfg["vendor"],
        "source": cfg["blueprint_url"],
        "disclaimer": CERTIPHY_DISCLAIMER,
        "mode": mode,
        "time_limit_minutes": time_limit_minutes if mode == "exam" else None,
        "questions": questions,
    }
    if CERTIPHY_CACHE_ENABLED and certiphy_set_cached:
        certiphy_set_cached(cache_key, response_obj)

    return response_obj


@app.post("/certiphy/exams/{exam_id}/generate-flashcards")
async def certiphy_generate_flashcards(exam_id: str, req: Dict[str, Any]):
    if not CERTIPHY_AVAILABLE or not certiphy_build_flashcards:
        raise HTTPException(status_code=503, detail=f"Certiphy API unavailable: {CERTIPHY_IMPORT_ERROR}")

    exam_id = exam_id.upper().strip()
    cfg = CERTIPHY_EXAMS.get(exam_id)
    if not cfg:
        raise HTTPException(status_code=404, detail="Unknown exam_id. Use GET /certiphy/exams")

    try:
        flashcards = await certiphy_build_flashcards(
            exam_id=exam_id,
            exam_name=cfg["name"],
            vendor=cfg["vendor"],
            questions=req.get("questions", []),
            use_ai=bool(req.get("use_ai", True)),
            max_cards=int(req.get("max_cards", 40)),
        )
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc))

    return {"exam_id": exam_id, "disclaimer": CERTIPHY_DISCLAIMER, "flashcards": flashcards}


if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)
